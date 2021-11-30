from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

num_cols = 7
num_rows = 6
offset_left = Cm(2.0)
offset_top = Cm(2.0)
diameter = Cm(1.0)
gap = Cm(2.0)
filename = 'calibration_pattern.pptx'
#type = "symmetrical"
#type = "asymmetrical"
type = "checker"

prs = Presentation();
SLIDE_BLANK = 6
slide_layout = prs.slide_layouts[SLIDE_BLANK]
slide = prs.slides.add_slide(slide_layout)

if type == "asymmetrical":
	for y in range(num_rows):
		for x in range(num_cols):
			if x % 2 == 1:
				shift = gap / 2
			else:
				shift = Cm(0)
			shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, offset_left + x * (gap/2), offset_top + y * (gap) + shift, diameter, diameter) 
			shape.fill.solid()
			shape.line.fill.background()
			shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
			shape.shadow.inherit = False
			shape.shadow.style = None
elif type == "checker":
	for y in range(num_rows):
		for x in range(num_cols):
			if (x+y) % 2 == 1:
				front_color = RGBColor(0xFF, 0xFF, 0xFF)
			else:
				front_color = RGBColor(0x00, 0x00, 0x00)
			shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, offset_left + x * diameter, offset_top + y * diameter, diameter, diameter) 
			shape.line.fill.background()
			shape.fill.solid()
			shape.fill.fore_color.rgb = front_color
			shape.shadow.inherit = False
			shape.shadow.style = None
else:
	for y in range(num_rows):
		for x in range(num_cols):
			shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, offset_left + x * (gap + diameter), offset_top + y * (gap + diameter), diameter, diameter) 
			shape.fill.solid()
			shape.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
			shape.shadow.inherit = False
			shape.shadow.style = None



prs.save(filename)
